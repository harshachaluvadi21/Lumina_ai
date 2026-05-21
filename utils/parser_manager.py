import re
import io
import os
import zipfile
import requests
import pdfplumber
import docx
from pptx import Presentation
from bs4 import BeautifulSoup
from youtube_transcript_api import YouTubeTranscriptApi
import google.generativeai as genai
import config

class ParserManager:
    @staticmethod
    def parse_pdf(file_bytes) -> str:
        """Parses PDF text page-by-page using pdfplumber for high-fidelity structural layouts."""
        text_content = []
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                page_text = page.extract_text()
                if page_text:
                    text_content.append(f"[Page {page_num}]\n{page_text}")
        return "\n\n".join(text_content)

    @staticmethod
    def parse_docx(file_bytes) -> str:
        """Parses Microsoft Word document paragraphs."""
        doc = docx.Document(io.BytesIO(file_bytes))
        text_content = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(text_content)

    @staticmethod
    def parse_pptx(file_bytes) -> str:
        """Parses Microsoft PowerPoint slides and speaker shapes."""
        prs = Presentation(io.BytesIO(file_bytes))
        text_content = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            slide_text.append(paragraph.text)
            if slide_text:
                text_content.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))
        return "\n\n".join(text_content)

    @staticmethod
    def parse_txt(file_bytes) -> str:
        """Parses generic plain text files."""
        return file_bytes.decode("utf-8", errors="ignore")

    @staticmethod
    def parse_web_url(url: str) -> str:
        """Scrapes text headings and paragraphs from raw website HTML."""
        headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"}
        res = requests.get(url, headers=headers, timeout=10)
        res.raise_for_status()
        
        soup = BeautifulSoup(res.text, "html.parser")
        
        # Remove script and style tags
        for element in soup(["script", "style", "nav", "footer", "header"]):
            element.decompose()
            
        title = soup.title.string if soup.title else "Web Article"
        
        paragraphs = []
        for tag in soup.find_all(['h1', 'h2', 'h3', 'h4', 'p']):
            text = tag.get_text().strip()
            if text:
                paragraphs.append(f"{tag.name.upper()}: {text}")
                
        if not paragraphs:
            # Fallback to general stripped text if no headers/paragraphs found
            return f"Title: {title}\n\n" + soup.get_text(separator="\n", strip=True)
            
        return f"Title: {title}\nURL: {url}\n\n" + "\n\n".join(paragraphs)

    @staticmethod
    def extract_youtube_video_id(url: str) -> str:
        """Extracts the YouTube 11-character video ID using standard regular expressions."""
        pattern = r'(?:https?://)?(?:www\.)?(?:youtube\.com/(?:[^/]+/.+/|(?:v|e(?:mbed)?)/|.*[?&]v=)|youtu\.be/)([^"&?/ ]{11})'
        match = re.match(pattern, url)
        if match:
            return match.group(1)
        return ""

    @classmethod
    def parse_youtube(cls, url: str) -> str:
        """Extracts full timestamped transcription from a YouTube URL."""
        video_id = cls.extract_youtube_video_id(url)
        if not video_id:
            raise ValueError("Invalid YouTube URL. Could not extract video ID.")
            
        try:
            api = YouTubeTranscriptApi()
            transcript_list = api.fetch(video_id)
            full_text = []
            for entry in transcript_list:
                start_sec = int(entry.start)
                minutes = start_sec // 60
                seconds = start_sec % 60
                timestamp = f"[{minutes:02d}:{seconds:02d}]"
                full_text.append(f"{timestamp} {entry.text}")
            return f"YouTube Transcript (ID: {video_id})\nURL: {url}\n\n" + "\n".join(full_text)
        except Exception as e:
            # Try to grab metadata if transcript api fails
            raise ValueError(f"Could not retrieve YouTube Transcript: {str(e)}. "
                             "Please check if the video has English subtitles/captions enabled.")

    @staticmethod
    def parse_zip_codebase(file_bytes) -> str:
        """Unpacks ZIP codebase archives, parses structural files, and maps contents."""
        allowed_extensions = {
            '.py', '.js', '.ts', '.jsx', '.tsx', '.html', '.css', '.json',
            '.java', '.c', '.cpp', '.h', '.cs', '.go', '.rs', '.sh', '.md',
            '.yml', '.yaml', '.properties', '.sql', '.gradle', '.xml'
        }
        ignored_folders = {'__pycache__', '.git', 'node_modules', '.venv', 'venv', 'env', 'dist', 'build'}
        
        extracted_text = []
        
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
            for file_info in z.infolist():
                if file_info.is_dir():
                    continue
                    
                parts = Path_parts = file_info.filename.split('/')
                # Check for ignored folders
                if any(ignored in parts for ignored in ignored_folders):
                    continue
                    
                _, ext = os.path.splitext(file_info.filename)
                if ext.lower() in allowed_extensions:
                    try:
                        with z.open(file_info) as f:
                            content = f.read().decode('utf-8', errors='ignore')
                            if content.strip():
                                extracted_text.append(
                                    f"=========================================\n"
                                    f"FILE: {file_info.filename}\n"
                                    f"=========================================\n"
                                    f"{content}\n"
                                )
                    except Exception:
                        continue
                        
        if not extracted_text:
            return "No readable source code files found in the ZIP archive."
            
        return "\n\n".join(extracted_text)

    @staticmethod
    def parse_audio_with_gemini(file_bytes, original_filename: str) -> str:
        """Transcribes uploaded audio files using Gemini's native multimodal capabilities."""
        if not config.is_gemini_available():
            raise ValueError("Gemini API Key is missing. Audio transcription cannot proceed.")
            
        # Determine extension and save a temporary file locally
        _, ext = os.path.splitext(original_filename)
        if not ext:
            ext = ".mp3"  # default fallback
            
        temp_path = os.path.join(config.UPLOAD_DIR, f"temp_audio_{os.getpid()}{ext}")
        with open(temp_path, "wb") as f:
            f.write(file_bytes)
            
        try:
            # Configure generative AI SDK
            genai.configure(api_key=config.GEMINI_API_KEY)
            
            # Upload file using Gemini File API
            print(f"Uploading {original_filename} to Gemini File API...")
            uploaded_file = genai.upload_file(path=temp_path)
            
            # Select working model name
            working_model = "gemini-1.5-flash"
            try:
                available = [m.name for m in genai.list_models() if "generateContent" in m.supported_generation_methods]
                for cand in available:
                    if cand == "models/gemini-1.5-flash" or cand == "gemini-1.5-flash":
                        working_model = cand
                        break
                else:
                    for cand in available:
                        if "gemini-1.5-flash" in cand:
                            working_model = cand
                            break
                    else:
                        for cand in available:
                            if "gemini-1.5" in cand:
                                working_model = cand
                                break
                        else:
                            if available:
                                working_model = available[0]
            except Exception as e:
                print(f"Error listing models for audio transcription: {e}")
                
            model = genai.GenerativeModel(working_model)
            prompt = (
                "Please transcribe this audio file completely, accurately, and word-for-word. "
                "Output ONLY the text transcript of the audio. Do not summarize, "
                "do not add commentary or pleasantries."
            )
            
            print("Requesting native audio transcription from Gemini...")
            response = model.generate_content([uploaded_file, prompt])
            
            # Clean up the file from the Gemini servers
            genai.delete_file(uploaded_file.name)
            
            return response.text if response.text else "Failed to extract readable audio text."
            
        finally:
            # Clean up local file
            if os.path.exists(temp_path):
                os.remove(temp_path)
